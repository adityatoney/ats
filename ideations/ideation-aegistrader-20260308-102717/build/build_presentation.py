"""
Build script for AegisTrader presentation.
Generates a polished PowerPoint from the ideation session vision document.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Color Palette ---
DARK_BG = RGBColor(0x0B, 0x1D, 0x3A)       # Deep navy
DARKER_BG = RGBColor(0x07, 0x14, 0x2A)      # Darker navy for variation
LIGHT_BG = RGBColor(0x10, 0x27, 0x4F)       # Slightly lighter navy
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xC4, 0xDE)     # Light steel blue
ACCENT_BLUE = RGBColor(0x4F, 0xA8, 0xF5)    # Bright blue accent
ACCENT_TEAL = RGBColor(0x00, 0xD4, 0xAA)    # Teal/green accent
ACCENT_PURPLE = RGBColor(0xA8, 0x7C, 0xF5)  # Purple accent
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x4F)  # Orange accent
ACCENT_GREEN = RGBColor(0x4F, 0xF5, 0x7A)   # Green accent
MUTED_TEXT = RGBColor(0x7B, 0x93, 0xB3)      # Muted text

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set slide background to a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape as a background element."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bullet_color=ACCENT_BLUE, font_name="Calibri",
                    spacing=1.4):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.line_spacing = Pt(font_size * spacing)
        p.space_after = Pt(4)

        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "\u25B8  "
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.name = font_name

        # Item text — handle bold markers
        if "**" in item:
            parts = item.split("**")
            for j, part in enumerate(parts):
                if not part:
                    continue
                run = p.add_run()
                run.text = part
                run.font.size = Pt(font_size)
                run.font.name = font_name
                run.font.color.rgb = color
                run.font.bold = (j % 2 == 1)
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.name = font_name
            run.font.color.rgb = color

    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    """Add a thin accent line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_pill_box(slide, left, top, width, height, text, bg_color, text_color=WHITE,
                 font_size=14, bold=True):
    """Add a rounded rectangle pill with text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.margin_left = Pt(8)
    shape.text_frame.margin_right = Pt(8)
    shape.text_frame.margin_top = Pt(4)
    shape.text_frame.margin_bottom = Pt(4)
    return shape


# ============================================================
# BUILD PRESENTATION
# ============================================================

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Use blank layout
blank_layout = prs.slide_layouts[6]


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

# Accent line at top
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(4), ACCENT_BLUE)

# Title
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
             "Project AegisTrader", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
             "Growing Trading Personalities", font_size=28, color=ACCENT_TEAL,
             bold=False, alignment=PP_ALIGN.CENTER)

# Accent line
add_accent_line(slide, Inches(5.5), Inches(4.2), Inches(2.3), ACCENT_TEAL)

# Tagline
add_text_box(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.6),
             "Choose a personality. Grow it through experimentation. Coach it to improve.",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Date
add_text_box(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.5),
             "Ideation Session  |  March 8, 2026", font_size=14, color=MUTED_TEXT,
             alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Vision
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
             "Vision", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

# Core thesis box
add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(2.0), LIGHT_BG)
add_text_box(slide, Inches(1.2), Inches(1.7), Inches(10.8), Inches(0.5),
             "Core Thesis", font_size=20, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(10.8), Inches(1.2),
             "AegisTrader is not a backtesting platform with AI features. It is a platform for "
             "growing, experimenting on, and coaching trading personalities. The backtest engine "
             "provides credibility. The soul provides the experience.",
             font_size=17, color=LIGHT_GRAY, line_spacing=1.4)

# Governing principle
add_shape_bg(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(1.4), LIGHT_BG)
add_text_box(slide, Inches(1.2), Inches(4.1), Inches(10.8), Inches(0.5),
             "Governing Principle", font_size=20, color=ACCENT_TEAL, bold=True)
add_text_box(slide, Inches(1.2), Inches(4.6), Inches(10.8), Inches(0.7),
             '"Soul powered by backtest." Every soul claim must trace back to deterministic, '
             'reproducible evidence. The engine provides reproducibility and trust. '
             'The soul provides meaning and engagement.',
             font_size=17, color=LIGHT_GRAY, line_spacing=1.4)

# New category callout
add_pill_box(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.7),
             "This is a new product category  --  zero market analogs exist",
             ACCENT_PURPLE, WHITE, font_size=16, bold=True)


# ============================================================
# SLIDE 3: The Product Loop
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
             "The Product Loop", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.2), ACCENT_BLUE)

# Three main boxes for Choose -> Grow -> Coach
box_y = Inches(2.2)
box_h = Inches(3.5)
box_w = Inches(3.2)
gap = Inches(0.5)
start_x = Inches(1.0)

# CHOOSE
add_shape_bg(slide, start_x, box_y, box_w, box_h, LIGHT_BG)
add_text_box(slide, start_x + Inches(0.3), box_y + Inches(0.3), box_w - Inches(0.6), Inches(0.5),
             "1. CHOOSE", font_size=22, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, start_x + Inches(0.3), box_y + Inches(0.9), box_w - Inches(0.6), Inches(0.4),
             "Soul as the Product", font_size=14, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_list(slide, start_x + Inches(0.3), box_y + Inches(1.4), box_w - Inches(0.6), Inches(1.8),
                ["Archetype picker", "Trading philosophy selection", "Seed soul + strategy"],
                font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT_BLUE)

# Arrow 1
add_text_box(slide, start_x + box_w + Inches(0.1), box_y + Inches(1.4), Inches(0.4), Inches(0.5),
             "\u25B6", font_size=28, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

# GROW
grow_x = start_x + box_w + gap
add_shape_bg(slide, grow_x, box_y, box_w, box_h, LIGHT_BG)
add_text_box(slide, grow_x + Inches(0.3), box_y + Inches(0.3), box_w - Inches(0.6), Inches(0.5),
             "2. GROW", font_size=22, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, grow_x + Inches(0.3), box_y + Inches(0.9), box_w - Inches(0.6), Inches(0.4),
             "Counterfactual Forking", font_size=14, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_list(slide, grow_x + Inches(0.3), box_y + Inches(1.4), box_w - Inches(0.6), Inches(1.8),
                ["Formative experiences", "Trauma tests", "Soul forking & comparison"],
                font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT_TEAL)

# Arrow 2
add_text_box(slide, grow_x + box_w + Inches(0.1), box_y + Inches(1.4), Inches(0.4), Inches(0.5),
             "\u25B6", font_size=28, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

# COACH
coach_x = grow_x + box_w + gap
add_shape_bg(slide, coach_x, box_y, box_w, box_h, LIGHT_BG)
add_text_box(slide, coach_x + Inches(0.3), box_y + Inches(0.3), box_w - Inches(0.6), Inches(0.5),
             "3. COACH", font_size=22, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, coach_x + Inches(0.3), box_y + Inches(0.9), box_w - Inches(0.6), Inches(0.4),
             "User as Coach", font_size=14, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_list(slide, coach_x + Inches(0.3), box_y + Inches(1.4), box_w - Inches(0.6), Inches(1.8),
                ["Review & reflect", "Guided coaching", "Direct soul surgery"],
                font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT_PURPLE)

# Loop arrow text
add_text_box(slide, Inches(2.5), Inches(6.1), Inches(8), Inches(0.5),
             "Repeat: Next Run  /  Next Fork  /  Next Coaching Cycle",
             font_size=16, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4: Move 1 - Soul as the Product
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
             "Move 1: Soul as the Product", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3.5), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             '"Grow trading personalities that learn from experience" -- not "backtesting platform with AI agents"',
             font_size=16, color=ACCENT_TEAL, bold=False)

# Soul-First Onboarding steps
steps = [
    ("Archetype Picker", "0-2 min", "Choose a trading philosophy: Mean Reversion, Trend Following,\nEvent-Driven, or Defensive Value. Each = strategy.md + seed soul.md.", ACCENT_BLUE),
    ("Formative Experience", "2-5 min", "System runs a compressed 3-6 month backtest through a dramatic\nhistorical period. UI emphasizes soul formation, not PnL.", ACCENT_TEAL),
    ("Soul Diff Hero Screen", "5-10 min", 'Before/after showing how the agent changed. Every belief\nevidence-linked: "Reduces size after consecutive losses -- from losing 12% in week 3."', ACCENT_PURPLE),
    ("Fork Prompt", "", '"Want to keep training? Or fork this soul and see what happens\nwith a different experience?" Bridges to counterfactual forking.', ACCENT_ORANGE),
]

y = Inches(2.2)
for label, time_label, desc, color in steps:
    # Colored left bar
    add_shape_bg(slide, Inches(0.8), y, Pt(4), Inches(1.05), color)
    add_text_box(slide, Inches(1.1), y + Inches(0.05), Inches(3.0), Inches(0.4),
                 label, font_size=18, color=color, bold=True)
    if time_label:
        add_text_box(slide, Inches(4.0), y + Inches(0.05), Inches(1.5), Inches(0.4),
                     time_label, font_size=12, color=MUTED_TEXT)
    add_text_box(slide, Inches(1.1), y + Inches(0.45), Inches(11.0), Inches(0.6),
                 desc, font_size=14, color=LIGHT_GRAY, line_spacing=1.3)
    y += Inches(1.25)


# ============================================================
# SLIDE 5: Move 2 - Counterfactual Soul Forking
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
             "Move 2: Counterfactual Soul Forking", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(4.0), ACCENT_TEAL)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             "Fork beliefs -- not just trades -- to run \"what if\" experiments on trading psychology",
             font_size=16, color=ACCENT_TEAL)

# Trauma Test MVP
add_shape_bg(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(3.8), LIGHT_BG)
add_pill_box(slide, Inches(1.0), Inches(2.4), Inches(1.2), Inches(0.4), "MVP", ACCENT_BLUE, WHITE, 12)
add_text_box(slide, Inches(2.4), Inches(2.35), Inches(3.5), Inches(0.5),
             "Trauma Test", font_size=22, color=WHITE, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(3.0), Inches(5.0), Inches(2.8),
                ["Clone same archetype, different eras",
                 "Bull-market agent: aggressive, naive about drawdowns",
                 "Bear-forged agent: cautious, misses rallies",
                 "Neither is \"right\" -- divergence is the lesson",
                 "Implementation: same strategy, different time ranges"],
                font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT_BLUE)

# Selective Amnesia v2
add_shape_bg(slide, Inches(6.8), Inches(2.2), Inches(5.5), Inches(3.8), LIGHT_BG)
add_pill_box(slide, Inches(7.0), Inches(2.4), Inches(1.0), Inches(0.4), "v2", ACCENT_PURPLE, WHITE, 12)
add_text_box(slide, Inches(8.2), Inches(2.35), Inches(3.8), Inches(0.5),
             "Selective Amnesia", font_size=22, color=WHITE, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(3.0), Inches(5.0), Inches(2.8),
                ['"What if my agent forgot its worst month?"',
                 "Remove a period from a mature soul's experience",
                 "Re-derive soul from parameterized experience set",
                 "Tests whether scars are wisdom or damage",
                 "Deep emotional resonance for traders"],
                font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT_PURPLE)

# Signature UX callout
add_shape_bg(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.7), RGBColor(0x15, 0x30, 0x55))
add_text_box(slide, Inches(1.2), Inches(6.35), Inches(10.8), Inches(0.6),
             'Signature UX: Side-by-side soul comparison -- "Bull-market you believes X. Crash-market you believes Y."',
             font_size=15, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6: Move 3 - User as Coach
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
             "Move 3: User as Coach", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3.0), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             "The user is a coach developing a trading personality, not an admin approving file changes",
             font_size=16, color=ACCENT_TEAL)

# Three coaching modes
modes = [
    ("Let it Learn", "Passive -- MVP", "Post-game film review. Review soul\ndiffs after runs. Accept or flag\nspecific beliefs for revision.", ACCENT_BLUE, "self_learned"),
    ("Guided Reflection", "Active -- v1.5", "Natural language coaching before\nsoul update commits. Shape the\nreflection, not the conclusion.", ACCENT_TEAL, "coach_guided"),
    ("Soul Surgery", "Direct -- MVP", "Power-user direct editing of\nsoul.md/soul.json. Edits tagged\nas coach_override.", ACCENT_PURPLE, "coach_override"),
]

mode_w = Inches(3.5)
mode_x = Inches(0.8)
for label, phase, desc, color, tag in modes:
    add_shape_bg(slide, mode_x, Inches(2.2), mode_w, Inches(2.8), LIGHT_BG)
    add_text_box(slide, mode_x + Inches(0.3), Inches(2.35), mode_w - Inches(0.6), Inches(0.5),
                 label, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, mode_x + Inches(0.3), Inches(2.8), mode_w - Inches(0.6), Inches(0.3),
                 phase, font_size=12, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, mode_x + Inches(0.3), Inches(3.2), mode_w - Inches(0.6), Inches(1.5),
                 desc, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.4)
    mode_x += mode_w + Inches(0.4)

# Soul Studio
add_shape_bg(slide, Inches(0.8), Inches(5.3), Inches(5.5), Inches(1.6), LIGHT_BG)
add_text_box(slide, Inches(1.1), Inches(5.4), Inches(5.0), Inches(0.4),
             "Soul Studio UI (v2)", font_size=18, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(1.1), Inches(5.8), Inches(5.0), Inches(1.0),
                ["Left: Soul Timeline (color-coded by source)",
                 "Center: Belief Cards with confidence + evidence",
                 "Right: Evidence Drawer for any selected belief"],
                font_size=13, color=LIGHT_GRAY, bullet_color=ACCENT_ORANGE)

# Retention Hook
add_shape_bg(slide, Inches(6.8), Inches(5.3), Inches(5.5), Inches(1.6), LIGHT_BG)
add_text_box(slide, Inches(7.1), Inches(5.4), Inches(5.0), Inches(0.4),
             'Retention Hook: "Am I a Good Coach?"', font_size=18, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(7.1), Inches(5.8), Inches(5.0), Inches(1.0),
                ["Track coached vs organic beliefs over time",
                 "Surface coaching effectiveness metrics",
                 "Per-belief attribution of user intuitions"],
                font_size=13, color=LIGHT_GRAY, bullet_color=ACCENT_GREEN)


# ============================================================
# SLIDE 7: How They Fit Together
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
             "How They Fit Together", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3.0), ACCENT_BLUE)

# Flow diagram as text blocks connected by arrows
flow_items = [
    ("Archetype Picker", "Move 1", ACCENT_BLUE, Inches(4.8)),
    ("Formative Experience", "Move 1", ACCENT_BLUE, Inches(4.8)),
    ("Soul Diff Screen", "Move 1", ACCENT_BLUE, Inches(3.0)),
    ("Fork + Trauma Test", "Move 2", ACCENT_TEAL, Inches(7.5)),
    ("Soul Comparison", "Move 2", ACCENT_TEAL, Inches(7.5)),
]

y = Inches(1.6)
for label, move, color, x_pos in flow_items[:3]:
    add_pill_box(slide, x_pos, y, Inches(3.2), Inches(0.55), f"{label}  ({move})", color, WHITE, 13)
    y += Inches(0.7)
    if label != "Soul Diff Screen":
        add_text_box(slide, Inches(5.8), y - Inches(0.15), Inches(1.0), Inches(0.3),
                     "\u25BC", font_size=16, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)
        y += Inches(0.15)

# Fork branches
fork_y = Inches(3.3)
add_text_box(slide, Inches(3.8), fork_y + Inches(0.3), Inches(1.0), Inches(0.3),
             "\u25BC", font_size=16, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.2), fork_y - Inches(0.1), Inches(1.0), Inches(0.3),
             "\u25B6", font_size=16, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)

# Move 2 branch
add_pill_box(slide, Inches(8.5), fork_y - Inches(0.25), Inches(3.2), Inches(0.55),
             "Fork + Trauma Test  (Move 2)", ACCENT_TEAL, WHITE, 13)
add_text_box(slide, Inches(9.5), fork_y + Inches(0.35), Inches(1.0), Inches(0.3),
             "\u25BC", font_size=16, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)
add_pill_box(slide, Inches(8.5), fork_y + Inches(0.55), Inches(3.2), Inches(0.55),
             "Soul Comparison  (Move 2)", ACCENT_TEAL, WHITE, 13)

# Move 3 coaching
coach_y = fork_y + Inches(1.0)
add_pill_box(slide, Inches(2.8), coach_y, Inches(3.2), Inches(0.55),
             "Coaching  (Move 3)", ACCENT_PURPLE, WHITE, 13)

add_text_box(slide, Inches(2.0), coach_y + Inches(0.65), Inches(2.2), Inches(0.35),
             "Let it Learn", font_size=13, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.8), coach_y + Inches(0.65), Inches(2.6), Inches(0.35),
             "Guided Reflection", font_size=13, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.8), coach_y + Inches(0.65), Inches(2.0), Inches(0.35),
             "Soul Surgery", font_size=13, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)

# Convergence screen callout
add_shape_bg(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.4), RGBColor(0x15, 0x30, 0x55))
add_text_box(slide, Inches(1.2), Inches(5.7), Inches(10.8), Inches(0.4),
             "The Convergence Screen: Side-by-Side Soul Comparison",
             font_size=20, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.2), Inches(6.15), Inches(10.8), Inches(0.7),
             "All three moves converge on one screen: see two souls shaped by different experiences, "
             "read their divergent beliefs with evidence links, and coach either one. "
             "This is the product's defining interaction.",
             font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.4)


# ============================================================
# SLIDE 8: Key Design Decisions
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
             "Key Design Decisions", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3.0), ACCENT_BLUE)

decisions = [
    ("Evidence Chain as Shared Infrastructure", ACCENT_BLUE,
     ["Every belief in soul.json links to specific trades, branches, and runs",
      "Structured evidence_refs[]: run_id, trade_ids, branch_comparisons, confidence, source_type",
      "Enables trustworthy soul diffs, meaningful comparisons, and coaching attribution",
      "MVP requirement: evidence chain schema must ship in Phase 1"]),
    ("Deterministic / Non-Deterministic Boundary", ACCENT_TEAL,
     ["Two-layer architecture preserved from Concept_1.md",
      "Deterministic layer: accepted decisions, account state, fills, branch outcomes",
      "Non-deterministic layer: proposals, annotations, hypotheses, summaries",
      "Coaching operates in non-deterministic layer; hard constraints flow through approval gate"]),
    ("Source Tagging From Day One", ACCENT_PURPLE,
     ["Every belief carries source: self_learned | coach_override | coach_guided | counterfactual_derived",
      "Enables: agent pushback, coaching effectiveness, counterfactual provenance",
      "Non-negotiable MVP decision -- retrofitting provenance is architecturally painful"]),
]

y = Inches(1.5)
for title, color, items in decisions:
    add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(1.75), LIGHT_BG)
    add_shape_bg(slide, Inches(0.8), y, Pt(4), Inches(1.75), color)
    add_text_box(slide, Inches(1.2), y + Inches(0.08), Inches(10.8), Inches(0.4),
                 title, font_size=18, color=color, bold=True)
    add_bullet_list(slide, Inches(1.2), y + Inches(0.45), Inches(10.8), Inches(1.2),
                    items, font_size=13, color=LIGHT_GRAY, bullet_color=color, spacing=1.3)
    y += Inches(1.95)


# ============================================================
# SLIDE 9: Competitive Positioning
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
             "Competitive Positioning", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3.3), ACCENT_BLUE)

# Four moat pillars
pillars = [
    ("Deterministic\nSimulation Truth", "Reproducible, auditable\nbacktests validated by\nNautilusTrader pattern", ACCENT_BLUE),
    ("Evolving\nTrading Souls", "Versioned, diffable,\nevidence-linked agent\ndoctrine", ACCENT_TEAL),
    ("Branch DAG\nas Biography", "Git-like history for\nboth trades and beliefs", ACCENT_PURPLE),
    ("Coaching\nRelationship", "User shapes agent\ndevelopment through\nreview and guidance", ACCENT_ORANGE),
]

px = Inches(0.8)
pw = Inches(2.7)
for title, desc, color in pillars:
    add_shape_bg(slide, px, Inches(1.6), pw, Inches(2.8), LIGHT_BG)
    add_shape_bg(slide, px, Inches(1.6), pw, Pt(4), color)
    add_text_box(slide, px + Inches(0.2), Inches(1.8), pw - Inches(0.4), Inches(0.8),
                 title, font_size=17, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, px + Inches(0.2), Inches(2.7), pw - Inches(0.4), Inches(1.2),
                 desc, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.4)
    px += pw + Inches(0.3)

# Zero analogs callout
add_pill_box(slide, Inches(0.8), Inches(4.7), Inches(5.5), Inches(0.5),
             "Zero market analogs for evolving agent doctrine", ACCENT_TEAL, WHITE, 14)
add_pill_box(slide, Inches(6.7), Inches(4.7), Inches(5.5), Inches(0.5),
             "Zero market analogs for branch DAG biography", ACCENT_PURPLE, WHITE, 14)

# Market size
add_shape_bg(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.4), LIGHT_BG)
add_text_box(slide, Inches(1.2), Inches(5.6), Inches(10.8), Inches(0.4),
             "Market Context", font_size=18, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(1.2), Inches(6.0), Inches(10.8), Inches(0.8),
                ["Algorithmic trading market: **$24B** (2025) projected to **$44.5B** (2030), **13.2% CAGR**",
                 "AegisTrader targets the research-oriented builder segment with a product no existing platform addresses",
                 'Positioning leads with what is unique: **"Grow trading personalities that learn from experience"**'],
                font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT_ORANGE)


# ============================================================
# SLIDE 10: Boundaries
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
             "Boundaries: What AegisTrader is NOT", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(4.5), RGBColor(0xF5, 0x5E, 0x5E))

boundaries = [
    ("Not a brokerage", "Does not execute real-money trades in initial versions", RGBColor(0xF5, 0x5E, 0x5E)),
    ("Not financial advice", "Clearly labeled as paper trading and research tooling", RGBColor(0xF5, 0x5E, 0x5E)),
    ("Not an opaque AI black box", "Every soul belief links to evidence. Every decision is inspectable.\nThe soul is interpretation, not truth.", RGBColor(0xF5, 0x5E, 0x5E)),
    ("The soul is not the source of truth", "The deterministic engine is. The soul is a learned interpretation\nof truth -- derived, versioned, diffable, and replay-linked.", RGBColor(0xF5, 0x5E, 0x5E)),
]

y = Inches(1.8)
for title, desc, color in boundaries:
    add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(1.15), LIGHT_BG)
    add_shape_bg(slide, Inches(0.8), y, Pt(4), Inches(1.15), color)
    add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(4.0), Inches(0.4),
                 title, font_size=18, color=color, bold=True)
    add_text_box(slide, Inches(1.3), y + Inches(0.5), Inches(10.5), Inches(0.6),
                 desc, font_size=15, color=LIGHT_GRAY, line_spacing=1.3)
    y += Inches(1.35)


# ============================================================
# SLIDE 11: Phasing
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
             "Phasing", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5), ACCENT_BLUE)

phases = [
    ("Phase 1", "Single-Agent Soul Evolution (MVP)", ACCENT_BLUE,
     ["4 archetype trading philosophies + formative experience backtest",
      "Soul diff hero screen with evidence links",
      "Trauma Test counterfactual + side-by-side soul comparison",
      "Passive coaching + Soul Surgery + source tagging"]),
    ("Phase 1.5", "Guided Coaching (NEW)", ACCENT_TEAL,
     ["Guided Reflection -- natural language coaching",
      "Challenge button on belief cards",
      "Enhanced Soul Studio layout"]),
    ("Phase 2", "Multi-Agent Competition", ACCENT_PURPLE,
     ["As defined in PRD + coaching effectiveness tracking",
      "Per-agent soul comparison across agents"]),
    ("Phase 3", "Live Paper Trading", ACCENT_ORANGE,
     ["Live soul updates with coaching review before commit",
      "Coaching effectiveness dashboard"]),
    ("Phase 4", "Governance & Advanced Features", MUTED_TEXT,
     ["Selective Amnesia counterfactual",
      "Full coaching effectiveness analytics",
      "Evaluate: Adversarial Soul Dynamics based on user demand"]),
]

y = Inches(1.5)
for phase_label, phase_title, color, items in phases:
    row_h = Inches(0.3 + len(items) * 0.32)
    add_shape_bg(slide, Inches(0.8), y, Inches(11.5), row_h, LIGHT_BG)
    add_shape_bg(slide, Inches(0.8), y, Pt(4), row_h, color)
    add_pill_box(slide, Inches(1.1), y + Inches(0.08), Inches(1.3), Inches(0.35), phase_label, color, WHITE, 12)
    add_text_box(slide, Inches(2.6), y + Inches(0.08), Inches(4.0), Inches(0.35),
                 phase_title, font_size=15, color=color, bold=True)
    add_bullet_list(slide, Inches(5.8), y + Inches(0.05), Inches(6.2), row_h - Inches(0.1),
                    items, font_size=12, color=LIGHT_GRAY, bullet_color=color, spacing=1.2)
    y += row_h + Inches(0.12)


# ============================================================
# SLIDE 12: Open Questions
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
             "Open Questions", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.2), ACCENT_ORANGE)

questions = [
    ("LLM Non-Determinism in Counterfactuals",
     "How many runs needed to establish stable soul divergences? What threshold for surfacing a difference as real vs noise?"),
    ("Guided Reflection Quality",
     "Can LLMs reliably integrate vague coaching (\"you're too cautious\") into coherent soul updates?"),
    ("Soul Timeline Noise Management",
     "After 20+ runs, the timeline will be dense. What summarization and grouping strategies keep it usable?"),
    ("Archetype Design",
     "Who designs seed strategy.md + soul.md for each archetype? These must be genuinely viable strategies."),
    ("Coaching Effectiveness Validity",
     "How much run history needed before coaching effectiveness metrics are statistically meaningful?"),
]

y = Inches(1.6)
for i, (title, desc) in enumerate(questions):
    add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(1.0), LIGHT_BG)
    add_text_box(slide, Inches(1.2), y + Inches(0.08), Inches(0.4), Inches(0.35),
                 str(i + 1), font_size=18, color=ACCENT_ORANGE, bold=True)
    add_text_box(slide, Inches(1.6), y + Inches(0.08), Inches(10.0), Inches(0.35),
                 title, font_size=16, color=ACCENT_ORANGE, bold=True)
    add_text_box(slide, Inches(1.6), y + Inches(0.45), Inches(10.0), Inches(0.5),
                 desc, font_size=14, color=LIGHT_GRAY, line_spacing=1.3)
    y += Inches(1.1)


# ============================================================
# SLIDE 13: Closing
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

# Accent line at top
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(4), ACCENT_TEAL)

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.0),
             "AegisTrader", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.5), Inches(3.0), Inches(2.3), ACCENT_TEAL)

# The three actions
actions_y = Inches(3.5)
action_items = [
    ("Choose", "a personality", ACCENT_BLUE),
    ("Grow", "it through experimentation", ACCENT_TEAL),
    ("Coach", "it to improve", ACCENT_PURPLE),
]

ax = Inches(1.5)
for verb, rest, color in action_items:
    add_text_box(slide, ax, actions_y, Inches(3.5), Inches(0.6),
                 verb, font_size=32, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, ax, actions_y + Inches(0.6), Inches(3.5), Inches(0.5),
                 rest, font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    ax += Inches(3.5)

# Tagline
add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.6),
             '"Soul powered by backtest."',
             font_size=22, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

# Session info
add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.4),
             "Ideation Session  |  March 8, 2026  |  Project AegisTrader",
             font_size=14, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
output_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
output_path = os.path.join(output_dir, "PRESENTATION_aegistrader.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
